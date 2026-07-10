from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageOps
import gc
import io
from typing import List, Tuple

# Paths
SCRIPT_DIR = Path(__file__).resolve().parent if "__file__" in globals() else Path(".").resolve()
ASSETS_DIR = SCRIPT_DIR / "assets"

# --- Image sizing ---------------------------------------------------------
# Slide images are placed into a fixed 24.89 x 12.87cm box. At 24.89cm
# (~9.8in), 1600px gives ~163 DPI on the long edge - comfortably above the
# 150 DPI print threshold. Anything larger is discarded by PowerPoint's
# renderer anyway, but python-pptx retains every original blob in memory
# until prs.save(), so oversized sources drive peak RSS on Render.
MAX_IMAGE_EDGE_PX = 1600
JPEG_QUALITY = 88

TEMPLATE_PATHS = {
    "new_campaign": ASSETS_DIR / "PoP Report - New Campaign Template.pptx",
    "artwork_update": ASSETS_DIR / "PoP Report - Artwork Update Template.pptx",
}
DEFAULT_REPORT_TYPE = "new_campaign"

LOGO_PATH = ASSETS_DIR / "GAWK LOGO (PURPLE).png"
BACKGROUND_PATH = ASSETS_DIR / "Background.jpg"

# Colours
PURPLE = RGBColor(0x54, 0x2D, 0x54)
GAWK_GREEN = RGBColor(0xD7, 0xDF, 0x23)


def convert_date(date_str: str) -> str:
    """Convert DDMMYY to 'MONTH YYYY' in uppercase."""
    try:
        return datetime.strptime(date_str, "%d%m%y").strftime("%B %Y").upper()
    except Exception:
        return "INVALID DATE"


def display_date(date_str: str) -> str:
    """Convert DDMMYY to DD/MM/YY for on-slide display."""
    try:
        return datetime.strptime(date_str, "%d%m%y").strftime("%d/%m/%y")
    except Exception:
        return "INVALID DATE"


def parse_filename(path_like) -> dict | None:
    """
    Parse filenames in the format:
    'Site Name - Site Code - Client - Campaign - DDMMYY - Type[ - OptionalSuffix]'.
    Only the first 6 parts are meaningful for the PoP logic.
    """
    # Accept either a Path or an object with a .name attribute (e.g. UploadedFile)
    if hasattr(path_like, "name"):
        name = Path(path_like.name)
    else:
        name = Path(path_like)

    parts = name.stem.split(" - ")
    if len(parts) < 6:
        # Too short to be valid
        return None

    try:
        site_name = f"{parts[0].strip()} - {parts[1].strip()}"
        client = parts[2].strip()
        campaign = parts[3].strip()
        live_date_raw = parts[4].strip()
        return {
            "site_name": site_name,
            "client": client,
            "campaign": campaign,
            "live_date": live_date_raw,
            "month_year": convert_date(live_date_raw),
            "live_date_display": display_date(live_date_raw),
        }
    except Exception:
        return None


def prepare_image(source, max_edge: int = MAX_IMAGE_EDGE_PX) -> io.BytesIO:
    """
    Load an image, correct EXIF orientation, downscale to fit `max_edge`,
    and return it as an in-memory JPEG buffer.

    python-pptx holds every image blob it is given until the presentation is
    saved, so the size of what we hand it - not the size of the file on the
    wire - determines peak memory. Downscaling here is what keeps long
    campaigns viable without capping the number of photos.

    Accepts anything Pillow can open: a Path, a file-like object, or a
    Streamlit UploadedFile.
    """
    img = None
    converted = None
    try:
        img = Image.open(source)
        # Honour EXIF rotation before we discard the metadata.
        img = ImageOps.exif_transpose(img)

        if max(img.size) > max_edge:
            img.thumbnail((max_edge, max_edge), Image.LANCZOS)

        # JPEG cannot store an alpha channel or a palette.
        if img.mode != "RGB":
            converted = img.convert("RGB")
            img.close()
            img = converted
            converted = None

        buffer = io.BytesIO()
        img.save(buffer, format="JPEG", quality=JPEG_QUALITY, optimize=True)
        buffer.seek(0)
        return buffer
    finally:
        if converted is not None:
            converted.close()
        if img is not None:
            img.close()


def _sort_key_from_name(filename: str):
    """
    Sort key, derived from a filename string rather than a Path:
    1. By actual live date (earliest first)
    2. Then by base filename (first 5 parts)
    3. Then by suffix priority: Cam (0) < Mock (1) < Others (2)
    """
    try:
        name = filename.lower()
        parts = Path(filename).stem.split(" - ")
        base_key = " - ".join(parts[:5]).lower()
        suffix = name.split(" - ")[-1]

        if "cam" in suffix:
            suffix_priority = 0
        elif "mock" in suffix:
            suffix_priority = 1
        else:
            suffix_priority = 2

        date = datetime.strptime(parts[4], "%d%m%y")
        return (date, base_key, suffix_priority)
    except Exception:
        return (datetime.min, filename.lower(), 9)


def extract_live_date_priority(path: Path):
    """
    Sort key matching desktop script:
    1. By actual live date (earliest first)
    2. Then by base filename (first 5 parts)
    3. Then by suffix priority: Cam (0) < Mock (1) < Others (2)
    """
    try:
        name = path.name.lower()
        parts = path.stem.split(" - ")
        base_key = " - ".join(parts[:5]).lower()
        suffix = name.split(" - ")[-1]

        if "cam" in suffix:
            suffix_priority = 0
        elif "mock" in suffix:
            suffix_priority = 1
        else:
            suffix_priority = 2

        date_str = parts[4]
        date = datetime.strptime(date_str, "%d%m%y")
        return (date, base_key, suffix_priority)
    except Exception:
        return (datetime.min, path.name.lower(), 9)


def _add_front_slide_content(prs: Presentation, first_info: dict) -> None:
    first_slide = prs.slides[0]
    for shape in first_slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text
        if "Client Name" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = first_info["client"]
            r.font.name = "Montserrat"
            r.font.size = Pt(60)
            r.font.bold = True
            r.font.color.rgb = RGBColor(255, 255, 255)
        elif "DATE" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = first_info["month_year"]
            r.font.name = "Montserrat"
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = RGBColor(255, 255, 255)
        elif "Campaign Name" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = first_info["campaign"]
            r.font.name = "Montserrat"
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = GAWK_GREEN


def _add_text(slide, x, y, w, h, text, size=23, bold=True, color=RGBColor(255, 255, 255)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Montserrat"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def resolve_template(report_type: str) -> Path:
    """
    Map a report type to its .pptx template, failing loudly on an unknown type
    or a missing file rather than silently falling back to the wrong deck.
    """
    try:
        template = TEMPLATE_PATHS[report_type]
    except KeyError:
        valid = ", ".join(sorted(TEMPLATE_PATHS))
        raise ValueError(f"Unknown report type {report_type!r}. Expected one of: {valid}.") from None

    if not template.exists():
        raise FileNotFoundError(f"Template missing for {report_type!r}: {template.name}")

    return template


def generate_presentation_bytes(
    images: List[Tuple[str, io.BytesIO]],
    report_type: str = DEFAULT_REPORT_TYPE,
) -> Tuple[bytes, str]:
    """
    Core PoP logic. Takes (filename, JPEG buffer) pairs - already downscaled -
    and returns the finished .pptx as bytes.

    `report_type` selects the base template. Both templates share identical
    placeholder text and slide geometry; only the cover headline differs.

    Nothing is written to disk. Buffers are released as soon as python-pptx has
    copied them into the package.
    """
    template_path = resolve_template(report_type)

    images = sorted(images, key=lambda pair: _sort_key_from_name(pair[0]))

    if not images:
        raise FileNotFoundError("No JPG, JPEG, or PNG files found for PoP generation.")

    # Cover slide content comes from the first image whose filename parses.
    # Unparseable names sort to the front (their date falls back to datetime.min),
    # so indexing images[0] blindly would fail on an otherwise valid batch.
    first_info = next(
        (info for info in (parse_filename(n) for n, _ in images) if info),
        None,
    )
    if not first_info:
        raise ValueError(
            "No image filename matches the required convention. "
            "Cannot determine client/campaign/date for the cover slide."
        )

    prs = Presentation(template_path)
    blank_layout = prs.slide_layouts[5]

    _add_front_slide_content(prs, first_info)

    for filename, buffer in images:
        details = parse_filename(filename)
        if not details:
            buffer.close()
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Background strip
        slide.shapes.add_picture(
            str(BACKGROUND_PATH),
            Cm(0),
            Cm(-0.01),
            width=Cm(29.7),
            height=Cm(21),
        )

        # Main PoP image, scaled to fit the 24.89 x 12.87cm box at (3.4, 4.45).
        # Read dimensions from the JPEG header without decoding pixel data.
        buffer.seek(0)
        with Image.open(buffer) as probe:
            iw, ih = probe.size

        img_aspect = iw / ih
        box_aspect = 24.89 / 12.87

        if img_aspect > box_aspect:
            new_width = Cm(24.89)
            new_height = Cm(24.89 / img_aspect)
        else:
            new_height = Cm(12.87)
            new_width = Cm(12.87 * img_aspect)

        buffer.seek(0)
        slide.shapes.add_picture(
            buffer,
            Cm(3.4),
            Cm(4.45),
            width=new_width,
            height=new_height,
        )
        # python-pptx has copied the bytes into its own package part by now.
        buffer.close()

        # Gawk logo top-right
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Cm(23.8),
            Cm(1.52),
            width=Cm(4.49),
            height=Cm(1.46),
        )

        # Gawk green vertical strip
        rect = slide.shapes.add_shape(
            1,  # MSO_SHAPE_RECTANGLE
            Cm(0),
            Cm(0),
            Cm(1.22),
            Cm(21),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = GAWK_GREEN
        rect.line.fill.background()

        # Vertical 'PROOF OF POSTING'
        tb = slide.shapes.add_textbox(Cm(-9.73), Cm(9.96), Cm(20.71), Cm(0.94))
        tb.rotation = 270
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "PROOF OF POSTING"
        r.font.name = "Montserrat"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = PURPLE

        # Site + Live Date labels/values
        _add_text(
            slide,
            Cm(3.06),
            Cm(2.5),
            Cm(3.09),
            Cm(1.24),
            "Site:",
            color=GAWK_GREEN,
        )
        _add_text(
            slide,
            Cm(5.36),
            Cm(2.5),
            Cm(13.25),
            Cm(1.24),
            details["site_name"],
        )
        _add_text(
            slide,
            Cm(3.06),
            Cm(18),
            Cm(4.76),
            Cm(1.24),
            "Live Date:",
            color=GAWK_GREEN,
        )
        _add_text(
            slide,
            Cm(7.79),
            Cm(18),
            Cm(12.35),
            Cm(1.24),
            details["live_date_display"],
        )

    # Remove example slide (index 1) and move 'Gotta love rectangles' to the end
    if len(prs.slides) >= 3:
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[1])
        final_slide = prs.slides._sldIdLst[1]
        prs.slides._sldIdLst.remove(final_slide)
        prs.slides._sldIdLst.append(final_slide)

    safe_client = first_info["client"].strip()
    safe_month = first_info["month_year"]
    # Distinguish the two report types so an artwork update doesn't overwrite
    # a new campaign deck for the same client and month.
    suffix = " - Artwork Update" if report_type == "artwork_update" else ""
    output_name = f"PoP Report - {safe_client} ({safe_month}){suffix}.pptx"

    bio = io.BytesIO()
    prs.save(bio)

    # prs.save() is the peak: python-pptx serialises every retained image blob
    # into the zip buffer. Release the presentation before we materialise bytes.
    del prs
    gc.collect()

    data = bio.getvalue()
    bio.close()
    gc.collect()

    return data, output_name


def generate_presentation_from_uploads(
    uploaded_files,
    report_type: str = DEFAULT_REPORT_TYPE,
) -> Tuple[bytes, str]:
    """
    Wrapper for Streamlit: takes a list of UploadedFile objects and returns
    (pptx_bytes, suggested_filename).

    Images are downscaled one at a time and never touch the filesystem.
    """
    # Fail before decoding a single image if the template is wrong or missing.
    resolve_template(report_type)

    images: list[tuple[str, io.BytesIO]] = []
    try:
        for uf in uploaded_files:
            if Path(uf.name).suffix.lower() not in (".jpg", ".jpeg", ".png"):
                continue
            uf.seek(0)
            images.append((uf.name, prepare_image(uf)))

        return generate_presentation_bytes(images, report_type=report_type)
    except Exception:
        # generate_presentation_bytes closes buffers as it consumes them; on a
        # failure part-way through, close whatever is still open.
        for _, buffer in images:
            if not buffer.closed:
                buffer.close()
        raise
    finally:
        gc.collect()